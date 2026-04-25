from __future__ import annotations

import base64
from datetime import datetime
import inspect
import json
import mimetypes
import re
import shutil
import subprocess
import tempfile
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Annotated, Any, Callable, Literal, cast, get_args, get_origin, get_type_hints
from urllib.parse import urlparse

import fitz
import pyzipper
import requests
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pdfminer.high_level import extract_text as extract_pdf_text
from PIL import Image
from pptx import Presentation

from support_ope_agents.config.models import AppConfig
from support_ope_agents.tools.infer_log_pattern import build_default_infer_log_pattern_tool
from support_ope_agents.util.ai_chat_util_bridge import (
    analyze_image_files as analyze_image_files_with_ai_chat_util,
    analyze_office_files as analyze_office_files_with_ai_chat_util,
    analyze_pdf_files as analyze_pdf_files_with_ai_chat_util,
)


ToolCallable = Callable[..., Any]


TEXT_FILE_SUFFIXES = {
    ".txt",
    ".log",
    ".md",
    ".csv",
    ".json",
    ".yaml",
    ".yml",
    ".xml",
    ".html",
    ".htm",
    ".py",
    ".js",
    ".ts",
    ".sql",
}


@dataclass(frozen=True, slots=True)
class BuiltinTool:
    name: str
    description: str
    handler: ToolCallable
    input_schema: dict[str, Any]


def _annotation_description(annotation: Any) -> tuple[Any, str | None]:
    if get_origin(annotation) is not Annotated:
        return annotation, None

    metadata = get_args(annotation)
    if not metadata:
        return annotation, None

    base_annotation = metadata[0]
    description: str | None = None
    for item in metadata[1:]:
        if isinstance(item, str) and item.strip():
            description = item.strip()
            break
        field_description = getattr(item, "description", None)
        if isinstance(field_description, str) and field_description.strip():
            description = field_description.strip()
            break
    return base_annotation, description


def _json_schema_for_annotation(annotation: Any) -> dict[str, Any]:
    annotation, description = _annotation_description(annotation)
    origin = get_origin(annotation)
    args = get_args(annotation)

    if annotation is Any:
        schema: dict[str, Any] = {}
    elif annotation in {str, Path}:
        schema = {"type": "string"}
    elif annotation is int:
        schema = {"type": "integer"}
    elif annotation is float:
        schema = {"type": "number"}
    elif annotation is bool:
        schema = {"type": "boolean"}
    elif annotation is datetime:
        schema = {"type": "string", "format": "date-time"}
    elif origin is Literal:
        literal_values = list(args)
        schema = {"enum": literal_values}
        if literal_values:
            first = literal_values[0]
            if isinstance(first, bool):
                schema["type"] = "boolean"
            elif isinstance(first, int) and not isinstance(first, bool):
                schema["type"] = "integer"
            elif isinstance(first, float):
                schema["type"] = "number"
            elif isinstance(first, str):
                schema["type"] = "string"
    elif origin is list:
        item_annotation = args[0] if args else Any
        schema = {"type": "array", "items": _json_schema_for_annotation(item_annotation)}
    elif origin in {dict}:
        schema = {"type": "object"}
    elif origin in {tuple}:
        item_annotation = args[0] if args else Any
        schema = {"type": "array", "items": _json_schema_for_annotation(item_annotation)}
    else:
        union_members = [member for member in args if member is not type(None)]
        has_null = len(union_members) != len(args)
        if args and union_members:
            variants = [_json_schema_for_annotation(member) for member in union_members]
            if has_null:
                variants.append({"type": "null"})
            schema = {"anyOf": variants}
        else:
            schema = {}

    if description:
        schema["description"] = description
    return schema


def _json_safe_default(value: Any) -> Any:
    if value is None or isinstance(value, (str, int, float, bool)):
        return value
    if isinstance(value, (list, dict)):
        return value
    return str(value)


def _build_input_schema(handler: ToolCallable) -> dict[str, Any]:
    try:
        signature = inspect.signature(handler)
    except (TypeError, ValueError):
        return {}

    try:
        type_hints = get_type_hints(handler, include_extras=True)
    except Exception:
        type_hints = {}

    properties: dict[str, Any] = {}
    required: list[str] = []
    for parameter in signature.parameters.values():
        if parameter.kind in {inspect.Parameter.VAR_POSITIONAL, inspect.Parameter.VAR_KEYWORD}:
            continue
        annotation = type_hints.get(parameter.name, Any)
        parameter_schema = _json_schema_for_annotation(annotation)
        if parameter.default is not inspect.Signature.empty:
            parameter_schema["default"] = _json_safe_default(parameter.default)
        else:
            required.append(parameter.name)
        properties[parameter.name] = parameter_schema

    doc = inspect.getdoc(handler) or ""
    schema: dict[str, Any] = {
        "type": "object",
        "properties": properties,
        "additionalProperties": False,
    }
    if required:
        schema["required"] = required
    if doc.strip():
        schema["description"] = doc.strip().splitlines()[0].strip()
    return schema


def _builtin_tool(name: str, description: str, handler: ToolCallable) -> BuiltinTool:
    schema = _build_input_schema(handler)
    resolved_description = description.strip() or str(schema.get("description") or name)
    return BuiltinTool(name=name, description=resolved_description, handler=handler, input_schema=schema)


def _ensure_existing_paths(paths: list[str]) -> list[Path]:
    resolved: list[Path] = []
    for raw_path in paths:
        path = Path(raw_path).expanduser().resolve()
        if not path.exists():
            raise FileNotFoundError(f"File was not found: {path}")
        resolved.append(path)
    return resolved


def _truncate_text(text: str, limit: int) -> str:
    if len(text) <= limit:
        return text
    return text[: limit - 16] + "\n...[truncated]"


def _extract_docx_text(path: Path) -> str:
    document = DocxDocument(str(path))
    lines = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    return "\n".join(lines)


def _extract_pptx_text(path: Path) -> str:
    presentation = Presentation(str(path))
    lines: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                lines.append(f"[slide {slide_index}] {text.strip()}")
    return "\n".join(lines)


def _extract_xlsx_text(path: Path) -> str:
    workbook = load_workbook(path, read_only=True, data_only=True)
    lines: list[str] = []
    for sheet in workbook.worksheets:
        lines.append(f"[sheet] {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            values = [str(value).strip() for value in row if value is not None and str(value).strip()]
            if values:
                lines.append("\t".join(values))
    workbook.close()
    return "\n".join(lines)


def _extract_text_from_path(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in TEXT_FILE_SUFFIXES:
        return path.read_text(encoding="utf-8", errors="ignore")
    if suffix == ".pdf":
        return extract_pdf_text(str(path))
    if suffix == ".docx":
        return _extract_docx_text(path)
    if suffix == ".pptx":
        return _extract_pptx_text(path)
    if suffix in {".xlsx", ".xlsm"}:
        return _extract_xlsx_text(path)
    raise ValueError(f"Text extraction is not supported for file type: {path.suffix or '<none>'}")


def _coerce_url(value: Any) -> tuple[str, dict[str, str] | None]:
    if isinstance(value, str):
        return value, None
    if isinstance(value, dict):
        url = value.get("url") or value.get("uri") or value.get("href")
        headers = value.get("headers")
        if isinstance(url, str):
            normalized_headers = None
            if isinstance(headers, dict):
                normalized_headers = {str(key): str(item) for key, item in headers.items()}
            return url, normalized_headers
    raise ValueError(f"Unsupported URL entry: {value!r}")


def _suffix_from_url(url: str, response: requests.Response) -> str:
    candidate = Path(urlparse(url).path).suffix
    if candidate:
        return candidate
    content_type = response.headers.get("content-type")
    guessed = mimetypes.guess_extension((content_type or "").split(";", 1)[0].strip())
    return guessed or ".bin"


def _download_urls(config: AppConfig, url_entries: list[Any]) -> tuple[tempfile.TemporaryDirectory[str], list[Path]]:
    tmpdir = tempfile.TemporaryDirectory()
    downloaded: list[Path] = []
    for index, entry in enumerate(url_entries, start=1):
        url, headers = _coerce_url(entry)
        response = requests.get(url, headers=headers, timeout=config.tools.download_timeout_seconds)
        response.raise_for_status()
        suffix = _suffix_from_url(url, response)
        target = Path(tmpdir.name) / f"download_{index:04d}{suffix}"
        target.write_bytes(response.content)
        downloaded.append(target)
    return tmpdir, downloaded


def _resolve_output_dir(output_dir: str | None) -> Path | None:
    if output_dir is None or not str(output_dir).strip():
        return None
    return Path(output_dir).expanduser().resolve()


def _resolve_libreoffice_command(config: AppConfig) -> str:
    candidates = [config.tools.libreoffice_command, "soffice", "libreoffice"]
    for candidate in candidates:
        if not candidate:
            continue
        resolved = shutil.which(candidate) or (candidate if Path(candidate).exists() else None)
        if resolved:
            return resolved
    raise RuntimeError(
        "LibreOffice command was not found. Set tools.libreoffice_command in config.yml or install soffice/libreoffice."
    )


def _collect_text_documents(config: AppConfig, paths: list[Path]) -> list[dict[str, Any]]:
    documents: list[dict[str, Any]] = []
    for path in paths:
        text = _truncate_text(_extract_text_from_path(path), config.tools.analysis_max_chars)
        documents.append(
            {
                "name": path.name,
                "path": str(path),
                "content": text,
                "content_length": len(text),
            }
        )
    return documents


def _detect_log_format_from_lines(lines: list[str]) -> dict[str, Any]:
    joined = "\n".join(lines)
    format_scores = {
        "syslog": sum(1 for line in lines if re.match(r"^[A-Z][a-z]{2}\s+\d{1,2}\s+\d{2}:\d{2}:\d{2}\s+", line)),
        "log4j": sum(1 for line in lines if re.match(r"^\d{4}-\d{2}-\d{2}[ T]\d{2}:\d{2}:\d{2}(?:[,.]\d{3,6})?\s+(?:TRACE|DEBUG|INFO|WARN|ERROR|FATAL)\b", line)),
        "iso8601": sum(1 for line in lines if re.match(r"^\d{4}-\d{2}-\d{2}T\d{2}:\d{2}:\d{2}", line)),
        "java_stacktrace": sum(1 for line in lines if re.match(r"^\s+at\s+[\w.$_]+\([^\n]+:\d+\)$", line)),
    }
    primary_format = max(format_scores, key=lambda item: format_scores[item]) if any(format_scores.values()) else "unknown"
    has_java_stacktrace = bool(re.search(r"^\s+at\s+[\w.$_]+\([^\n]+:\d+\)$", joined, flags=re.MULTILINE))

    timestamp_regex = r"\d{4}-\d{2}-\d{2}[ T]\d{2}:\d{2}:\d{2}(?:[,.]\d{3,6})?"
    if primary_format == "syslog":
        timestamp_regex = r"[A-Z][a-z]{2}\s+\d{1,2}\s+\d{2}:\d{2}:\d{2}"
    elif primary_format == "iso8601":
        timestamp_regex = r"\d{4}-\d{2}-\d{2}T\d{2}:\d{2}:\d{2}(?:\.\d+)?(?:Z|[+-]\d{2}:?\d{2})?"

    return {
        "primary_format": primary_format,
        "format_scores": format_scores,
        "has_java_stacktrace": has_java_stacktrace,
        "generated_patterns": {
            "timestamp": timestamp_regex,
            "severity": r"\b(?:TRACE|DEBUG|INFO|WARN|ERROR|FATAL)\b",
            "java_exception": r"\b[\w.$]+(?:Exception|Error)\b",
            "java_stack_frame": r"^\s+at\s+[\w.$_]+\([^\n]+:\d+\)$",
        },
    }


def _parse_timestamp_value(raw_value: str, time_format: str | None = None) -> datetime | None:
    value = raw_value.strip()
    if not value:
        return None
    if time_format:
        try:
            return datetime.strptime(value, time_format)
        except ValueError:
            pass

    normalized = value.replace(",", ".")
    try:
        return datetime.fromisoformat(normalized)
    except ValueError:
        pass

    candidates = (
        "%Y-%m-%d %H:%M:%S.%f",
        "%Y-%m-%d %H:%M:%S",
        "%Y-%m-%d %H:%M",
        "%Y/%m/%d %H:%M:%S.%f",
        "%Y/%m/%d %H:%M:%S",
        "%Y/%m/%d %H:%M",
        "%Y-%m-%dT%H:%M:%S.%f",
        "%Y-%m-%dT%H:%M:%S",
        "%Y%m%d %H:%M:%S",
    )
    for candidate in candidates:
        try:
            return datetime.strptime(normalized, candidate)
        except ValueError:
            continue
    return None


def _sanitize_output_component(value: str) -> str:
    collapsed = re.sub(r"[^0-9A-Za-z._-]+", "-", value.strip())
    sanitized = collapsed.strip("-._")
    return sanitized or "value"


def _extract_log_records_in_time_range(
    text: str,
    *,
    header_pattern: str,
    timestamp_start: int,
    timestamp_end: int,
    range_start: str,
    range_end: str,
    time_format: str | None = None,
) -> dict[str, Any]:
    if timestamp_start < 0 or timestamp_end <= timestamp_start:
        raise ValueError("timestamp_start and timestamp_end must define a valid non-empty slice")

    start_time = _parse_timestamp_value(range_start, time_format)
    end_time = _parse_timestamp_value(range_end, time_format)
    if start_time is None or end_time is None:
        raise ValueError("range_start and range_end must be parseable timestamps")
    if start_time > end_time:
        raise ValueError("range_start must be earlier than or equal to range_end")

    compiled = re.compile(header_pattern)
    lines = text.splitlines()
    records: list[dict[str, Any]] = []
    current_record: dict[str, Any] | None = None
    unmatched_preamble: list[str] = []

    for line_number, line in enumerate(lines, start=1):
        if compiled.search(line):
            if current_record is not None:
                records.append(current_record)
            current_record = {
                "start_line": line_number,
                "lines": [line],
                "header_line": line,
            }
            continue
        if current_record is not None:
            cast(list[str], current_record["lines"]).append(line)
        else:
            unmatched_preamble.append(line)

    if current_record is not None:
        records.append(current_record)

    extracted_records: list[dict[str, Any]] = []
    parse_failures: list[dict[str, Any]] = []
    for record in records:
        header_line = str(record["header_line"])
        timestamp_text = header_line[timestamp_start:timestamp_end]
        parsed_timestamp = _parse_timestamp_value(timestamp_text, time_format)
        if parsed_timestamp is None:
            parse_failures.append(
                {
                    "start_line": record["start_line"],
                    "timestamp_text": timestamp_text,
                    "header_line": header_line,
                }
            )
            continue
        if start_time <= parsed_timestamp <= end_time:
            extracted_records.append(
                {
                    "start_line": record["start_line"],
                    "end_line": record["start_line"] + len(cast(list[str], record["lines"])) - 1,
                    "timestamp": parsed_timestamp.isoformat(),
                    "text": "\n".join(cast(list[str], record["lines"])),
                }
            )

    return {
        "total_lines": len(lines),
        "matched_record_count": len(extracted_records),
        "record_count": len(records),
        "parse_failure_count": len(parse_failures),
        "parse_failures": parse_failures[:10],
        "records": extracted_records,
        "unmatched_preamble_line_count": len(unmatched_preamble),
        "time_range": {
            "start": start_time.isoformat(),
            "end": end_time.isoformat(),
        },
    }


def _search_log_with_patterns(
    text: str,
    *,
    generated_patterns: dict[str, str],
    search_terms: list[str] | None,
    match_limit: int,
) -> dict[str, list[dict[str, Any]]]:
    results: dict[str, list[dict[str, Any]]] = {}
    lines = text.splitlines()
    named_patterns = {
        "severity": generated_patterns["severity"],
        "java_exception": generated_patterns["java_exception"],
    }
    if search_terms:
        named_patterns["search_terms"] = "|".join(re.escape(term) for term in search_terms if term)

    for name, pattern in named_patterns.items():
        if not pattern:
            continue
        compiled = re.compile(pattern)
        matches: list[dict[str, Any]] = []
        for line_number, line in enumerate(lines, start=1):
            if compiled.search(line):
                matches.append({"line_number": line_number, "line": line})
                if len(matches) >= match_limit:
                    break
        results[name] = matches
    return results


def build_builtin_tools(config: AppConfig) -> dict[str, BuiltinTool]:
    infer_log_pattern = build_default_infer_log_pattern_tool(config)

    async def analyze_image_files(
        file_list: list[str],
        prompt: str,
        detail: str = "auto",
    ) -> str:
        paths = _ensure_existing_paths(file_list)
        return await analyze_image_files_with_ai_chat_util(config, [str(path) for path in paths], prompt, detail)

    async def analyze_pdf_files(
        pdf_path_list: list[str],
        prompt: str,
        detail: str = "auto",
    ) -> str:
        paths = _ensure_existing_paths(pdf_path_list)
        return await analyze_pdf_files_with_ai_chat_util(config, [str(path) for path in paths], prompt, detail)

    async def analyze_office_files(
        office_path_list: list[str],
        prompt: str,
        detail: str = "auto",
    ) -> str:
        paths = _ensure_existing_paths(office_path_list)
        return await analyze_office_files_with_ai_chat_util(config, [str(path) for path in paths], prompt, detail)

    async def analyze_image_urls(
        image_path_urls: list[Any],
        prompt: str,
        detail: str = "auto",
    ) -> str:
        tmpdir, paths = _download_urls(config, image_path_urls)
        try:
            return await analyze_image_files_with_ai_chat_util(config, [str(path) for path in paths], prompt, detail)
        finally:
            tmpdir.cleanup()

    async def analyze_pdf_urls(
        pdf_path_urls: list[Any],
        prompt: str,
        detail: str = "auto",
    ) -> str:
        tmpdir, paths = _download_urls(config, pdf_path_urls)
        try:
            return await analyze_pdf_files_with_ai_chat_util(config, [str(path) for path in paths], prompt, detail)
        finally:
            tmpdir.cleanup()

    async def analyze_office_urls(
        office_path_urls: list[Any],
        prompt: str,
        detail: str = "auto",
    ) -> str:
        tmpdir, paths = _download_urls(config, office_path_urls)
        try:
            return await analyze_office_files_with_ai_chat_util(config, [str(path) for path in paths], prompt, detail)
        finally:
            tmpdir.cleanup()

    async def extract_text_from_file(file_path: str) -> str:
        path = _ensure_existing_paths([file_path])[0]
        return _extract_text_from_path(path)

    async def extract_base64_to_text(extension: str, base64_data: str) -> str:
        suffix = extension if extension.startswith(".") else f".{extension}"
        raw = base64.b64decode(base64_data)
        with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as handle:
            temp_path = Path(handle.name)
            handle.write(raw)
        try:
            return _extract_text_from_path(temp_path)
        finally:
            temp_path.unlink(missing_ok=True)

    async def list_zip_contents(file_path: str) -> list[str]:
        path = _ensure_existing_paths([file_path])[0]
        with zipfile.ZipFile(path) as archive:
            return archive.namelist()

    async def extract_zip(file_path: str, extract_to: str, password: str | None = None) -> bool:
        archive_path = _ensure_existing_paths([file_path])[0]
        target_dir = Path(extract_to).expanduser().resolve()
        target_dir.mkdir(parents=True, exist_ok=True)
        if password:
            with pyzipper.AESZipFile(archive_path) as archive:
                archive.extractall(target_dir, pwd=password.encode("utf-8"))
            return True
        with zipfile.ZipFile(archive_path) as archive:
            archive.extractall(target_dir)
        return True

    async def create_zip(file_paths: list[str], output_zip: str, password: str | None = None) -> bool:
        paths = _ensure_existing_paths(file_paths)
        output_path = Path(output_zip).expanduser().resolve()
        output_path.parent.mkdir(parents=True, exist_ok=True)

        if password:
            with pyzipper.AESZipFile(output_path, "w", compression=zipfile.ZIP_DEFLATED, encryption=pyzipper.WZ_AES) as archive:
                archive.setpassword(password.encode("utf-8"))
                for path in paths:
                    if path.is_dir():
                        for child in path.rglob("*"):
                            if child.is_file():
                                archive.write(child, arcname=str(child.relative_to(path.parent)))
                    else:
                        archive.write(path, arcname=path.name)
            return True

        with zipfile.ZipFile(output_path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
            for path in paths:
                if path.is_dir():
                    for child in path.rglob("*"):
                        if child.is_file():
                            archive.write(child, arcname=str(child.relative_to(path.parent)))
                else:
                    archive.write(path, arcname=path.name)
        return True

    async def convert_office_files_to_pdf(
        office_path_list: list[str],
        output_dir: str | None = None,
        dry_run: bool = False,
    ) -> list[dict[str, str]]:
        paths = _ensure_existing_paths(office_path_list)
        target_root = _resolve_output_dir(output_dir)
        planned: list[dict[str, str]] = []
        for path in paths:
            pdf_dir = target_root or path.parent
            pdf_path = pdf_dir / f"{path.stem}.pdf"
            planned.append({"source_path": str(path), "pdf_path": str(pdf_path)})
        if dry_run:
            return planned

        command = _resolve_libreoffice_command(config)
        for path in paths:
            pdf_dir = target_root or path.parent
            pdf_dir.mkdir(parents=True, exist_ok=True)
            completed = subprocess.run(
                [command, "--headless", "--convert-to", "pdf", "--outdir", str(pdf_dir), str(path)],
                capture_output=True,
                text=True,
                check=False,
            )
            if completed.returncode != 0:
                raise RuntimeError(completed.stderr.strip() or completed.stdout.strip() or f"Failed to convert {path}")
        return planned

    async def convert_pdf_files_to_images(
        pdf_path_list: list[str],
        output_dir: str | None = None,
        dry_run: bool = False,
        dpi: int = 144,
    ) -> list[dict[str, Any]]:
        paths = _ensure_existing_paths(pdf_path_list)
        output_root = _resolve_output_dir(output_dir)
        results: list[dict[str, Any]] = []
        scale = dpi / 72.0
        matrix = fitz.Matrix(scale, scale)
        for path in paths:
            image_dir = (output_root / f"{path.stem}_pages") if output_root is not None else (path.parent / f"{path.stem}_pages")
            with fitz.open(path) as document:
                page_total = len(document)
                image_paths = [str(image_dir / f"{path.stem}_page_{index:04d}.png") for index in range(1, page_total + 1)]
                if not dry_run:
                    image_dir.mkdir(parents=True, exist_ok=True)
                    for index in range(1, page_total + 1):
                        page = document.load_page(index - 1)
                        pixmap = page.get_pixmap(matrix=matrix)
                        pixmap.save(str(image_dir / f"{path.stem}_page_{index:04d}.png"))
            results.append(
                {
                    "source_path": str(path),
                    "image_dir": str(image_dir),
                    "image_paths": image_paths,
                }
            )
        return results

    async def detect_log_format_and_search(
        file_path: Annotated[str, "解析対象のログファイルパス"],
        search_terms: Annotated[list[str] | None, "追加で検索するキーワード一覧"] = None,
        sample_line_limit: Annotated[int, "ログ先頭から形式判定に使う最大行数"] = 100,
        match_limit: Annotated[int, "パターンごとに返す最大ヒット件数"] = 50,
    ) -> str:
        """Detect log format and search matching records."""
        path = _ensure_existing_paths([file_path])[0]
        text = _extract_text_from_path(path)
        lines = text.splitlines()
        sample_lines = lines[:sample_line_limit]
        detection = _detect_log_format_from_lines(sample_lines)
        search_results = _search_log_with_patterns(
            text,
            generated_patterns=detection["generated_patterns"],
            search_terms=search_terms,
            match_limit=match_limit,
        )
        result = {
            "file_path": str(path),
            "sample_line_limit": sample_line_limit,
            "sample_preview": sample_lines[:10],
            "detected_format": detection["primary_format"],
            "format_scores": detection["format_scores"],
            "has_java_stacktrace": detection["has_java_stacktrace"],
            "generated_patterns": detection["generated_patterns"],
            "search_terms": search_terms or [],
            "search_results": search_results,
        }
        return json.dumps(result, ensure_ascii=False)

    async def infer_log_header_pattern(
        file_path: Annotated[str, "解析対象のログファイルパス"],
        sample_line_limit: Annotated[int, "ログ先頭から推定に使う最大行数"] = 100,
    ) -> str:
        """Infer a log header pattern and timestamp slice from sample lines."""
        return infer_log_pattern(file_path=file_path, sample_line_limit=sample_line_limit)

    async def extract_log_time_range(
        file_path: Annotated[str, "抽出対象のログファイルパス"],
        workspace_path: Annotated[str, "派生ファイルを書き出す workspace ルート"],
        header_pattern: Annotated[str, "各ログレコード先頭行に一致する正規表現"],
        timestamp_start: Annotated[int, "先頭行内の時刻文字列の開始位置"],
        timestamp_end: Annotated[int, "先頭行内の時刻文字列の終了位置"],
        range_start: Annotated[str, "抽出開始時刻"],
        range_end: Annotated[str, "抽出終了時刻"],
        time_format: Annotated[str | None, "datetime.strptime 互換の時刻書式"] = None,
        output_subdir: Annotated[str, "抽出ログを書き出す artifacts 配下サブディレクトリ"] = "log_extracts",
        output_filename: Annotated[str | None, "抽出結果の出力ファイル名"] = None,
    ) -> str:
        """Extract log records whose timestamps fall within the requested range."""
        source_path = _ensure_existing_paths([file_path])[0]
        workspace_root = Path(workspace_path).expanduser().resolve()
        text = _extract_text_from_path(source_path)
        extraction = _extract_log_records_in_time_range(
            text,
            header_pattern=header_pattern,
            timestamp_start=timestamp_start,
            timestamp_end=timestamp_end,
            range_start=range_start,
            range_end=range_end,
            time_format=time_format,
        )

        artifacts_dir = workspace_root / config.data_paths.artifacts_subdir / output_subdir.strip("/")
        artifacts_dir.mkdir(parents=True, exist_ok=True)
        derived_filename = output_filename or (
            f"{source_path.stem}_{_sanitize_output_component(range_start)}_{_sanitize_output_component(range_end)}{source_path.suffix or '.log'}"
        )
        output_path = artifacts_dir / derived_filename
        rendered = "\n\n".join(str(record["text"]) for record in extraction["records"])
        output_path.write_text(rendered + ("\n" if rendered else ""), encoding="utf-8")

        result = {
            "status": "matched" if extraction["matched_record_count"] else "unavailable",
            "file_path": str(source_path),
            "workspace_path": str(workspace_root),
            "output_path": str(output_path),
            "header_pattern": header_pattern,
            "timestamp_start": timestamp_start,
            "timestamp_end": timestamp_end,
            "time_format": time_format,
            "range_start": range_start,
            "range_end": range_end,
            **extraction,
        }
        return json.dumps(result, ensure_ascii=False)

    return {
        "analyze_image_files": _builtin_tool("analyze_image_files", "Analyze local image files", analyze_image_files),
        "analyze_pdf_files": _builtin_tool("analyze_pdf_files", "Analyze local PDF files", analyze_pdf_files),
        "analyze_office_files": _builtin_tool("analyze_office_files", "Analyze local Office files", analyze_office_files),
        "convert_office_files_to_pdf": _builtin_tool(
            "convert_office_files_to_pdf",
            "Convert Office files to PDF",
            convert_office_files_to_pdf,
        ),
        "convert_pdf_files_to_images": _builtin_tool(
            "convert_pdf_files_to_images",
            "Convert PDF files to page images",
            convert_pdf_files_to_images,
        ),
        "analyze_image_urls": _builtin_tool("analyze_image_urls", "Analyze image URLs", analyze_image_urls),
        "analyze_pdf_urls": _builtin_tool("analyze_pdf_urls", "Analyze PDF URLs", analyze_pdf_urls),
        "analyze_office_urls": _builtin_tool("analyze_office_urls", "Analyze Office URLs", analyze_office_urls),
        "extract_text_from_file": _builtin_tool("extract_text_from_file", "Extract text from a local file", extract_text_from_file),
        "extract_base64_to_text": _builtin_tool(
            "extract_base64_to_text",
            "Extract text from base64-encoded file content",
            extract_base64_to_text,
        ),
        "list_zip_contents": _builtin_tool("list_zip_contents", "List ZIP archive contents", list_zip_contents),
        "extract_zip": _builtin_tool("extract_zip", "Extract ZIP archive", extract_zip),
        "create_zip": _builtin_tool("create_zip", "Create ZIP archive", create_zip),
        "detect_log_format_and_search": _builtin_tool(
            "detect_log_format_and_search",
            "Detect log format from the first lines, generate regex patterns, and search the log",
            detect_log_format_and_search,
        ),
        "infer_log_header_pattern": _builtin_tool(
            "infer_log_header_pattern",
            "Infer a log header regex and timestamp slice from the first lines of a log file",
            infer_log_header_pattern,
        ),
        "extract_log_time_range": _builtin_tool(
            "extract_log_time_range",
            "Extract records in a timestamp range from a log file and save them into the workspace artifacts",
            extract_log_time_range,
        ),
    }