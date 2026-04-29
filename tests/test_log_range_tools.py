from __future__ import annotations

import asyncio
import base64
from datetime import datetime
import json
import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch
import zipfile

from ai_chat_util.ai_chat_util_base.ai_chat_util_models import ChatResponse
from docx import Document as DocxDocument
from openpyxl import Workbook
from pptx import Presentation

from support_desk_agent.agents.production.investigate_agent import InvestigateAgent, InvestigateAgentTools
from support_desk_agent.config.models import AppConfig
from support_desk_agent.tools.builtin_tools import build_builtin_tools
from support_desk_agent.tools.infer_log_pattern import build_default_infer_log_pattern_tool
from support_desk_agent.tools.registry import ToolRegistry
from support_desk_agent.util.ai_chat_util_bridge import build_ai_chat_util_config
from support_desk_agent.util.log_time_range import derive_log_extract_range_from_timeframe


class _FakeStructuredModel:
    def invoke(self, _messages):
        return {
            "header_pattern": r"^\d+\s+\[[^\]]+\]\s+(?:TRACE|DEBUG|INFO|WARN|ERROR|FATAL)\s+\d{4}-\d{2}-\d{2}T",
            "timestamp_start": 16,
            "timestamp_end": 39,
            "timestamp_format": "%Y-%m-%dT%H:%M:%S.%f",
            "confidence": 0.93,
            "reason": "先頭行で時刻が固定位置にあります。",
        }


class _FakeModel:
    def with_structured_output(self, _schema):
        return _FakeStructuredModel()


class _FakeTimeRangeStructuredModel:
    def invoke(self, _messages):
        return {
            "range_start": "2026-04-20T18:00:00",
            "range_end": "2026-04-20T21:00:00",
            "reason": "昨日の夕方を 18:00-21:00 と解釈しました。",
        }


class _FakeTimeRangeModel:
    def with_structured_output(self, _schema):
        return _FakeTimeRangeStructuredModel()


class LogRangeToolTests(unittest.TestCase):
    def _build_config(self) -> AppConfig:
        return AppConfig.model_validate(
            {
                "llm": {"provider": "openai", "model": "gpt-4.1", "api_key": "sk-test-value"},
                "config_paths": {},
                "data_paths": {},
                "interfaces": {},
                "agents": {},
            }
        )

    def test_extract_log_time_range_saves_matching_records(self) -> None:
        config = self._build_config()
        tools = build_builtin_tools(config)
        handler = tools["extract_log_time_range"].handler

        with tempfile.TemporaryDirectory() as tmpdir:
            workspace_path = Path(tmpdir)
            log_path = workspace_path / "vdp.log"
            log_path.write_text(
                "100 [main] INFO 2025-10-21T20:32:35.845 start\n"
                "startup detail\n"
                "200 [main] ERROR 2025-10-21T20:55:12.313 failure\n"
                "com.example.CacheException: boom\n"
                "300 [main] INFO 2025-10-21T21:00:00.000 end\n",
                encoding="utf-8",
            )

            raw = asyncio.run(
                handler(
                    str(log_path),
                    str(workspace_path),
                    r"^\d+\s+\[[^\]]+\]\s+(?:TRACE|DEBUG|INFO|WARN|ERROR|FATAL)\s+\d{4}-\d{2}-\d{2}T",
                    16,
                    39,
                    "2025-10-21T20:50:00.000",
                    "2025-10-21T20:56:00.000",
                    "%Y-%m-%dT%H:%M:%S.%f",
                )
            )
            payload = json.loads(raw)

            self.assertEqual(payload["status"], "matched")
            self.assertEqual(payload["matched_record_count"], 1)
            output_path = Path(payload["output_path"])
            self.assertTrue(output_path.exists())
            rendered = output_path.read_text(encoding="utf-8")
            self.assertIn("2025-10-21T20:55:12.313 failure", rendered)
            self.assertIn("com.example.CacheException: boom", rendered)

    def test_extract_log_time_range_accepts_iso_range_without_matching_time_format(self) -> None:
        config = self._build_config()
        tools = build_builtin_tools(config)
        handler = tools["extract_log_time_range"].handler

        with tempfile.TemporaryDirectory() as tmpdir:
            workspace_path = Path(tmpdir)
            log_path = workspace_path / "vdp.log"
            log_path.write_text(
                "100 [main] INFO 2025-10-21T20:32:35.845 start\n"
                "200 [main] ERROR 2025-10-21T20:55:12.313 failure\n",
                encoding="utf-8",
            )

            raw = asyncio.run(
                handler(
                    str(log_path),
                    str(workspace_path),
                    r"^\d+\s+\[[^\]]+\]\s+(?:TRACE|DEBUG|INFO|WARN|ERROR|FATAL)\s+\d{4}-\d{2}-\d{2}T",
                    16,
                    39,
                    "2025-10-21T20:50:00",
                    "2025-10-21T20:56:00",
                    "%Y-%m-%dT%H:%M:%S.%f",
                )
            )
            payload = json.loads(raw)

            self.assertEqual(payload["matched_record_count"], 1)

    def test_infer_log_header_pattern_returns_structured_payload(self) -> None:
        config = self._build_config()
        tool = build_default_infer_log_pattern_tool(config)

        with tempfile.TemporaryDirectory() as tmpdir:
            log_path = Path(tmpdir) / "vdp.log"
            log_path.write_text(
                "100 [main] INFO 2025-10-21T20:32:35.845 start\n",
                encoding="utf-8",
            )

            with patch("support_desk_agent.tools.infer_log_pattern.build_chat_openai_model", return_value=_FakeModel()):
                payload = json.loads(tool(file_path=str(log_path), sample_line_limit=20))

        self.assertEqual(payload["status"], "matched")
        self.assertEqual(payload["timestamp_start"], 16)
        self.assertEqual(payload["timestamp_end"], 39)
        self.assertEqual(payload["timestamp_format"], "%Y-%m-%dT%H:%M:%S.%f")
        self.assertTrue(payload["header_pattern"])

    def test_builtin_tools_expose_input_schema_from_annotations(self) -> None:
        config = self._build_config()
        tools = build_builtin_tools(config)

        infer_schema = tools["infer_log_header_pattern"].input_schema
        self.assertEqual(infer_schema["type"], "object")
        self.assertIn("file_path", infer_schema["properties"])
        self.assertEqual(infer_schema["properties"]["file_path"]["type"], "string")
        self.assertEqual(infer_schema["properties"]["file_path"]["description"], "解析対象のログファイルパス")
        self.assertIn("file_path", infer_schema["required"])

        extract_schema = tools["extract_log_time_range"].input_schema
        self.assertEqual(extract_schema["properties"]["timestamp_start"]["type"], "integer")
        self.assertEqual(extract_schema["properties"]["output_subdir"]["default"], "log_extracts")

    def test_registry_exposes_new_investigate_tools(self) -> None:
        config = self._build_config()
        registry = ToolRegistry(config)

        tools = {tool.name: tool for tool in registry.get_tools("InvestigateAgent")}

        self.assertIn("list_zip_contents", tools)
        self.assertIn("extract_zip", tools)
        self.assertIn("create_zip", tools)
        self.assertIn("analyze_image_files", tools)
        self.assertIn("analyze_pdf_files", tools)
        self.assertIn("analyze_office_files", tools)
        self.assertIn("convert_office_files_to_pdf", tools)
        self.assertIn("convert_pdf_files_to_images", tools)
        self.assertIn("detect_log_format_and_search", tools)
        self.assertIn("infer_log_header_pattern", tools)
        self.assertIn("extract_log_time_range", tools)
        self.assertIn("file_path", tools["infer_log_header_pattern"].input_schema["properties"])

    def test_ai_chat_util_bridge_uses_support_config_as_primary_source(self) -> None:
        config = self._build_config()

        bridged = build_ai_chat_util_config(config)

        self.assertEqual(bridged.llm.provider, config.llm.provider)
        self.assertEqual(bridged.llm.completion_model, config.llm.model)
        self.assertEqual(bridged.llm.api_key, config.llm.api_key)
        self.assertEqual(bridged.office2pdf.libreoffice_path, config.tools.libreoffice_command)

    def test_analyze_pdf_files_delegates_to_ai_chat_util(self) -> None:
        config = self._build_config()
        tools = build_builtin_tools(config)
        handler = tools["analyze_pdf_files"].handler

        with tempfile.TemporaryDirectory() as tmpdir:
            pdf_path = Path(tmpdir) / "sample.pdf"
            pdf_path.write_bytes(b"%PDF-1.4\n")

            with patch(
                "support_desk_agent.tools.builtin_tools.analyze_pdf_files_with_ai_chat_util",
                return_value="delegated via ai-chat-util",
            ) as delegated:
                result = asyncio.run(handler([str(pdf_path)], "check this", "auto"))

        self.assertEqual(result, "delegated via ai-chat-util")
        delegated.assert_called_once()

    def test_chat_response_text_normalization_uses_output_property(self) -> None:
        response = ChatResponse.model_validate({"output": " normalized text "})

        self.assertEqual(response.output, " normalized text ")

    def test_list_zip_contents_delegates_to_ai_chat_util(self) -> None:
        config = self._build_config()
        tools = build_builtin_tools(config)
        handler = tools["list_zip_contents"].handler

        with tempfile.TemporaryDirectory() as tmpdir:
            archive_path = Path(tmpdir) / "sample.zip"
            with zipfile.ZipFile(archive_path, "w") as archive:
                archive.writestr("a.txt", "hello")

            with patch(
                "support_desk_agent.tools.builtin_tools.list_zip_contents_with_ai_chat_util",
                return_value=["a.txt"],
            ) as delegated:
                result = asyncio.run(handler(str(archive_path)))

        self.assertEqual(result, ["a.txt"])
        delegated.assert_called_once_with(str(archive_path.resolve()))

    def test_create_zip_delegates_to_ai_chat_util(self) -> None:
        config = self._build_config()
        tools = build_builtin_tools(config)
        handler = tools["create_zip"].handler

        with tempfile.TemporaryDirectory() as tmpdir:
            source_path = Path(tmpdir) / "input.txt"
            source_path.write_text("hello", encoding="utf-8")
            output_path = Path(tmpdir) / "nested" / "archive.zip"

            with patch(
                "support_desk_agent.tools.builtin_tools.create_zip_with_ai_chat_util",
                return_value=True,
            ) as delegated:
                result = asyncio.run(handler([str(source_path)], str(output_path), "secret"))

        self.assertTrue(result)
        delegated.assert_called_once_with([str(source_path.resolve())], str(output_path.resolve()), "secret")

    def test_extract_text_from_file_docx_keeps_compat_output(self) -> None:
        config = self._build_config()
        handler = build_builtin_tools(config)["extract_text_from_file"].handler

        with tempfile.TemporaryDirectory() as tmpdir:
            docx_path = Path(tmpdir) / "sample.docx"
            document = DocxDocument()
            document.add_paragraph("first line")
            document.add_paragraph("")
            document.add_paragraph("second line")
            document.save(docx_path)

            result = asyncio.run(handler(str(docx_path)))

        self.assertEqual(result, "first line\nsecond line")

    def test_extract_text_from_file_pptx_keeps_compat_output(self) -> None:
        config = self._build_config()
        handler = build_builtin_tools(config)["extract_text_from_file"].handler

        with tempfile.TemporaryDirectory() as tmpdir:
            pptx_path = Path(tmpdir) / "sample.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            slide.shapes.title.text = "overview"
            textbox = slide.shapes.add_textbox(0, 0, 100, 100)
            textbox.text_frame.text = "details"
            presentation.save(pptx_path)

            result = asyncio.run(handler(str(pptx_path)))

        self.assertEqual(result, "[slide 1] overview\n[slide 1] details")

    def test_extract_text_from_file_xlsx_keeps_compat_output(self) -> None:
        config = self._build_config()
        handler = build_builtin_tools(config)["extract_text_from_file"].handler

        with tempfile.TemporaryDirectory() as tmpdir:
            xlsx_path = Path(tmpdir) / "sample.xlsx"
            workbook = Workbook()
            sheet = workbook.active
            sheet.title = "Summary"
            sheet.append(["col1", "col2"])
            sheet.append(["a", None])
            workbook.save(xlsx_path)

            result = asyncio.run(handler(str(xlsx_path)))

        self.assertEqual(result, "[sheet] Summary\ncol1\tcol2\na")

    def test_extract_base64_to_text_keeps_compat_output(self) -> None:
        config = self._build_config()
        handler = build_builtin_tools(config)["extract_base64_to_text"].handler

        payload = base64.b64encode("hello\nworld\n".encode("utf-8")).decode("ascii")

        result = asyncio.run(handler(".txt", payload))

        self.assertEqual(result, "hello\nworld\n")

    def test_convert_office_files_to_pdf_delegates_to_ai_chat_util(self) -> None:
        config = self._build_config()
        tools = build_builtin_tools(config)
        handler = tools["convert_office_files_to_pdf"].handler

        with tempfile.TemporaryDirectory() as tmpdir:
            source_path = Path(tmpdir) / "input.docx"
            source_path.write_text("dummy", encoding="utf-8")

            with patch(
                "support_desk_agent.tools.builtin_tools.convert_office_files_to_pdf_with_ai_chat_util",
                return_value=[{"source_path": str(source_path.resolve()), "pdf_path": str(source_path.with_suffix('.pdf').resolve())}],
            ) as delegated:
                result = asyncio.run(handler([str(source_path)], None, True))

        self.assertEqual(len(result), 1)
        delegated.assert_called_once_with(config, [str(source_path.resolve())], None, True)

    def test_detect_log_format_and_search_delegates_to_ai_chat_util(self) -> None:
        config = self._build_config()
        tools = build_builtin_tools(config)
        handler = tools["detect_log_format_and_search"].handler

        with tempfile.TemporaryDirectory() as tmpdir:
            log_path = Path(tmpdir) / "app.log"
            log_path.write_text("INFO start\n", encoding="utf-8")

            with patch(
                "support_desk_agent.tools.builtin_tools.detect_log_format_and_search_with_ai_chat_util",
                return_value=json.dumps({"detected_format": "log4j"}, ensure_ascii=False),
            ) as delegated:
                result = asyncio.run(handler(str(log_path), ["ERROR"], 10, 5))

        self.assertEqual(json.loads(result)["detected_format"], "log4j")
        delegated.assert_called_once_with(config, str(log_path.resolve()), ["ERROR"], 10, 5)

    def test_infer_log_header_pattern_delegates_to_ai_chat_util(self) -> None:
        config = self._build_config()
        tools = build_builtin_tools(config)
        handler = tools["infer_log_header_pattern"].handler

        with tempfile.TemporaryDirectory() as tmpdir:
            log_path = Path(tmpdir) / "app.log"
            log_path.write_text("INFO start\n", encoding="utf-8")

            with patch(
                "support_desk_agent.tools.builtin_tools.infer_log_header_pattern_with_ai_chat_util",
                return_value=json.dumps({"status": "matched", "header_pattern": "^INFO", "timestamp_start": 0, "timestamp_end": 4}, ensure_ascii=False),
            ) as delegated:
                result = asyncio.run(handler(str(log_path), 20))

        self.assertEqual(json.loads(result)["status"], "matched")
        delegated.assert_called_once_with(config, str(log_path.resolve()), 20)

    def test_investigate_agent_appends_extraction_summary(self) -> None:
        config = self._build_config()
        agent = InvestigateAgent(
            config=config,
            tools=InvestigateAgentTools(
                detect_log_format_tool=lambda *_args, **_kwargs: json.dumps(
                    {
                        "detected_format": "unknown",
                        "search_results": {"severity": [], "java_exception": []},
                    },
                    ensure_ascii=False,
                ),
                infer_log_header_pattern_tool=lambda *_args, **_kwargs: json.dumps(
                    {
                        "header_pattern": r"^\d+\s+\[[^\]]+\]\s+(?:TRACE|DEBUG|INFO|WARN|ERROR|FATAL)\s+\d{4}-\d{2}-\d{2}T",
                        "timestamp_start": 16,
                        "timestamp_end": 39,
                        "timestamp_format": "%Y-%m-%dT%H:%M:%S.%f",
                    },
                    ensure_ascii=False,
                ),
                extract_log_time_range_tool=lambda *_args, **_kwargs: json.dumps(
                    {
                        "output_path": "/tmp/extracted.log",
                        "matched_record_count": 2,
                    },
                    ensure_ascii=False,
                ),
            ),
        )

        with tempfile.TemporaryDirectory() as tmpdir:
            workspace_path = Path(tmpdir)
            evidence_dir = workspace_path / ".evidence"
            evidence_dir.mkdir()
            (evidence_dir / "vdp.log").write_text("sample\n", encoding="utf-8")

            result = agent.execute(
                {
                    "workspace_path": str(workspace_path),
                    "raw_issue": "vdp.log のエラーを見て",
                    "log_extract_range_start": "2025-10-21T20:50:00.000",
                    "log_extract_range_end": "2025-10-21T20:56:00.000",
                }
            )

        summary = str(result.get("log_analysis_summary") or "")
        self.assertIn("/tmp/extracted.log", summary)
        self.assertIn("指定時間帯", summary)

    def test_derive_log_extract_range_uses_llm_for_ambiguous_timeframe(self) -> None:
        config = self._build_config()
        with patch("support_desk_agent.util.log_time_range.build_chat_openai_model", return_value=_FakeTimeRangeModel()):
            derived = derive_log_extract_range_from_timeframe(
                "昨日の夕方に発生しました。",
                config=config,
                reference_datetime=datetime(2026, 4, 21, 9, 0, 0),
            )

        self.assertEqual(derived, ("2026-04-20T18:00:00", "2026-04-20T21:00:00"))


if __name__ == "__main__":
    unittest.main()